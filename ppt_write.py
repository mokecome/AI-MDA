from pptx import Presentation
from pptx.util import Inches
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


# 創建一份新的簡報物件
prs = Presentation()

# 背景與數據圖片路徑
folder_path='ppt'
front_page = f"{folder_path}/Front_page.png"
bg_image_path = f"{folder_path}/PPT_BG.jpg"


# 函數來設置投影片背景圖片
def set_slide_background(slide, image_path):
    left = top = Inches(0)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    # 添加背景圖片
    bg = slide.shapes.add_picture(image_path, left, top, slide_width, slide_height)
    # 將背景圖片移動到最底層
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)


#  主封面頁-----------------------------1111111--------------
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]

# 設置背景圖片
set_slide_background(slide, front_page)


# 目錄頁--------------------------------222222 ---------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "內容"
content = slide.placeholders[1]
content.text = "B、C、B+C肝個案追蹤人數與分析\n收案男女比例\n個案追蹤管理分析\nDashboard\n常見問題問答\n結論"

# 設置背景圖片
set_slide_background(slide, bg_image_path)

# 設置字體大小和顏色
title_format = title.text_frame.paragraphs[0].font
title_format.size = Pt(36)
title_format.color.rgb = RGBColor(0, 0, 0)  # 黑色字體

for paragraph in content.text_frame.paragraphs:
    content_format = paragraph.font
    content_format.size = Pt(36)
    content_format.color.rgb = RGBColor(0, 0, 0)  # 黑色字體


# B、C、B+C肝個案追蹤人數與分析 ------------3333333----------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "B、C、B+C肝個案追蹤人數與分析"
content = slide.placeholders[1]

# 設置背景圖片
set_slide_background(slide, bg_image_path)

# 設置字體大小和顏色
title_format = title.text_frame.paragraphs[0].font
title_format.size = Pt(36)
title_format.color.rgb = RGBColor(0, 0, 0)  # 黑色字體

# 添加B_C_BC肝人數表.png圖片到第4頁並縮小10%
img_path = f"{folder_path}/B_C_BC.png"  # 圖片路徑
height = Inches(7.5)  # 原高度
width = Inches(11.25)  # 原寬度

# 縮小10%
new_height = height * 0.75
new_width = width * 0.75

# 計算新的位置使圖片置中
slide_width = prs.slide_width
slide_height = prs.slide_height
left = (slide_width - new_width) / 2
top = (slide_height - new_height) / 2 + Inches(0.5)  # 向下移動0.5英寸

# 添加圖片
slide.shapes.add_picture(img_path, left, top, new_width, new_height)

# 收案男女比例------------------4444444---------------

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "收案男女比例"


# 設置背景圖片
set_slide_background(slide, bg_image_path)

# 設置字體大小和顏色
title_format = title.text_frame.paragraphs[0].font
title_format.size = Pt(36)
title_format.color.rgb = RGBColor(0, 0, 0)  # 黑色字體


# 收案男女比例.png圖片到第5頁並縮小10%
img_path = f"./{folder_path}/fm_ratio.png"  # 圖片路徑
height = Inches(6.6)  # 原高度
width = Inches(11.5)  # 原寬度

# 縮小10%
new_height = height *  0.75
new_width = width *  0.75

# 計算新的位置使圖片置中
slide_width = prs.slide_width
slide_height = prs.slide_height
left = (slide_width - new_width) / 2
top = (slide_height - new_height) / 2 + Inches(0.5)  # 向下移動0.5英寸

# 添加圖片
slide.shapes.add_picture(img_path, left, top, new_width, new_height)

# 個案追蹤管理分析----- -----------5555555-----------------------

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "個案追蹤管理分析"


# 設置背景圖片
set_slide_background(slide, bg_image_path)

# 設置字體大小和顏色
title_format = title.text_frame.paragraphs[0].font
title_format.size = Pt(36)
title_format.color.rgb = RGBColor(0, 0, 0)  # 黑色字體

# 收案男女比例.png圖片到第5頁並縮小10%
img_path = f"./{folder_path}/manage.png"  # 圖片路徑
height = Inches(7.2)  # 原高度
width = Inches(11.5)  # 原寬度

# 縮小10%
new_height = height * 0.75
new_width = width * 0.75

# 計算新的位置使圖片置中
slide_width = prs.slide_width
slide_height = prs.slide_height
left = (slide_width - new_width) / 2
top = (slide_height - new_height) / 2 + Inches(0.5)  # 向下移動0.5英寸

# 添加圖片
slide.shapes.add_picture(img_path, left, top, new_width, new_height)

# Dashboard-----------------------6666666------------------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Dashboard"


# 設置背景圖片
set_slide_background(slide, bg_image_path)

# 設置字體大小和顏色
title_format = title.text_frame.paragraphs[0].font
title_format.size = Pt(36)


# 收案男女比例.png圖片到第5頁並縮小10%
img_path = f"./{folder_path}/Dashboard.png"  # 圖片路徑
height = Inches(4.7)  # 原高度
width = Inches(17.5)  # 原寬度

# 縮小10%
new_height = height * 0.55
new_width = width * 0.55

# 計算新的位置使圖片置中
slide_width = prs.slide_width
slide_height = prs.slide_height
left = (slide_width - new_width) / 2
top = (slide_height - new_height) / 2 + Inches(0.5)  # 向下移動0.5英寸

# 添加圖片
slide.shapes.add_picture(img_path, left, top, new_width, new_height)

# 常見問題問答 主頁--------------------77777---------------
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]

# 設置背景圖片
set_slide_background(slide, bg_image_path)

# 設置字體大小和顏色
title.text = "常見問題問答"
for paragraph in title.text_frame.paragraphs:
    title_format = paragraph.font
    title_format.size = Pt(42)
    # title_format.color.rgb = RGBColor(0, 0, 0)  # 黑色字體

# 常見問題問答 1-----------------888888-------------------

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "常見問題問答"
content = slide.placeholders[1]
content.text = "Q：超過一年以上未回診比例(以360天計算)"

# 新增回答的文本框
left = Inches(0.95)  # 左邊距離
top = Inches(2.5)   # 上邊距離
width = Inches(6) # 寬度
height = Inches(5) # 高度

textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
text_frame.text = "A：超過一年以上未回診比例 28%"

# 設置字體大小
for paragraph in text_frame.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(32)  # 設置字體大小為 32 點

# 設置背景圖片
set_slide_background(slide, bg_image_path)

# 設置字體大小和顏色
title_format = title.text_frame.paragraphs[0].font
title_format.size = Pt(36)
title_format.color.rgb = RGBColor(0, 0, 0)  # 黑色字體



# 常見問題問答 2 ------------9999999-----------------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "常見問題問答"
content = slide.placeholders[1]
content.text = "Q：追蹤總個案數中，未回診率"


# 新增回答的文本框
left = Inches(0.95)  # 左邊距離
top = Inches(2.5)   # 上邊距離
width = Inches(6) # 寬度
height = Inches(5) # 高度

textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
text_frame.text = "A：未回診率為 20%"

# 設置字體大小
for paragraph in text_frame.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(32)  # 設置字體大小為 32 點

# 設置背景圖片
set_slide_background(slide, bg_image_path)

# 設置字體大小和顏色
title_format = title.text_frame.paragraphs[0].font
title_format.size = Pt(36)
title_format.color.rgb = RGBColor(0, 0, 0)  # 黑色字體

# 常見問題問答 3 ------------10+1----------------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "常見問題問答"
content = slide.placeholders[1]
content.text = "Q：輕、中、重度脂肪肝的個案人數"


# 新增回答的文本框
left = Inches(0.95)  # 左邊距離
top = Inches(2.5)   # 上邊距離
width = Inches(6) # 寬度
height = Inches(5) # 高度

textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
text_frame.text = """A：\n輕度脂肪肝（fatty liver mild）：22例\n中度脂肪肝（fatty liver moderate）：168例\n重度脂肪肝（fatty liver severe）：74例"""


# 設置字體大小
for paragraph in text_frame.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(32)  # 設置字體大小為 32 點

# 設置背景圖片
set_slide_background(slide, bg_image_path)

# 設置字體大小和顏色
title_format = title.text_frame.paragraphs[0].font
title_format.size = Pt(36)
title_format.color.rgb = RGBColor(0, 0, 0)  # 黑色字體


# 結論  ------------11+1-----------------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "結論"
content = slide.placeholders[1]
content.text = """B肝患者的數量在過去幾年中有顯著波動，目前呈下降趨勢，但總體患者數量仍在增加。
大部分患者仍在追蹤管理中，結案和轉出的比例相對較低，表明需要進一步加強管理和治療的跟進。
在性別分佈上，B肝男性患者數量明顯多於女性，而C肝女性患者數量略高於男性。"""


# 設置背景圖片
set_slide_background(slide, bg_image_path)

# 設置字體大小和顏色
title_format = title.text_frame.paragraphs[0].font
title_format.size = Pt(36)
title_format.color.rgb = RGBColor(0, 0, 0)  # 黑色字體



# 感謝頁  ------------------------------------------------------------------------------------------------

slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]

# 設置背景圖片
set_slide_background(slide, bg_image_path)

# 設置字體大小和顏色
title.text = "以上報告，謝謝"
for paragraph in title.text_frame.paragraphs:
    title_format = paragraph.font
    title_format.size = Pt(42)
    title_format.color.rgb = RGBColor(0, 0, 0)  # 黑色字體

# 將簡報物件存檔----------------------------------------------------------------------------------------------------------------
prs.save("本月肝炎個管成效.pptx")
